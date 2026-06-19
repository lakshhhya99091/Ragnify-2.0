"""
Document Parser — supports PDF (with TABLE extraction), DOCX, JPG/PNG (OCR), PPTX
Specially enhanced for banking & tender documents:
  - Extracts tables as structured markdown-like text
  - Preserves numbered lists, quantities, manpower specs
  - Detects document type (tender, policy, report, etc.)
Returns: {
  "text_by_source": [(label, text), ...],
  "links": [url, ...],
  "doc_type": str,
  "tables": [(label, table_text), ...]
}
"""
import re
import logging
from pathlib import Path
from typing import List, Tuple, Dict

logger = logging.getLogger(__name__)

# ─────────────────────────────────────────────────────────────────────────────
# Document-type detection keywords
# ─────────────────────────────────────────────────────────────────────────────
_DOC_TYPE_PATTERNS = {
    "tender":  r"tender|RFP|RFQ|NIT|notice inviting|bid|bidder|procurement|BOQ|bill of quantities|rate contract|work order|scope of work|eligibility criteria",
    "policy":  r"policy|guidelines|rules|regulations|circular|directive|compliance|norm",
    "report":  r"annual report|balance sheet|profit.loss|financial statement|audit report|quarterly",
    "loan":    r"loan agreement|mortgage|EMI|interest rate|repayment|collateral|borrower|lender",
    "account": r"account statement|transaction|debit|credit|balance|IFSC|MICR|passbook",
    "kyc":     r"KYC|know your customer|identity|verification|PAN|Aadhaar|passport",
}


def _detect_doc_type(full_text: str) -> str:
    text_lower = full_text.lower()
    for dtype, pattern in _DOC_TYPE_PATTERNS.items():
        if re.search(pattern, text_lower, re.IGNORECASE):
            return dtype
    return "general"


# ─────────────────────────────────────────────────────────────────────────────
# Table formatter — converts PyMuPDF table to readable markdown-like text
# ─────────────────────────────────────────────────────────────────────────────
def _format_table(table_data: List[List], source: str) -> str:
    """Convert a 2D list (table rows) to clean text for embedding."""
    if not table_data:
        return ""

    lines = []
    # Use first row as header if it looks like one
    has_header = table_data and any(
        cell and str(cell).strip().isupper() or
        (cell and len(str(cell).strip()) < 40 and not any(c.isdigit() for c in str(cell)))
        for cell in (table_data[0] or [])
    )

    rows = table_data
    for i, row in enumerate(rows):
        # Clean cells
        cells = [str(c).strip() if c is not None else "" for c in row]
        # Skip completely empty rows
        if not any(cells):
            continue

        if i == 0 and has_header:
            lines.append("| " + " | ".join(cells) + " |")
            lines.append("|" + "|".join(["---"] * len(cells)) + "|")
        else:
            lines.append("| " + " | ".join(cells) + " |")

    return "\n".join(lines)


# ─────────────────────────────────────────────────────────────────────────────
# PDF Parser — enhanced with table extraction
# ─────────────────────────────────────────────────────────────────────────────
def parse_pdf(file_path: str) -> Dict:
    """Extract text (per page), tables, and hyperlinks from a PDF."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise RuntimeError("PyMuPDF not installed: pip install PyMuPDF")

    text_by_source: List[Tuple[str, str]] = []
    tables_by_source: List[Tuple[str, str]] = []
    links: List[str] = []
    all_text_parts = []

    doc = fitz.open(file_path)

    for page_num, page in enumerate(doc, start=1):
        source_label = f"Page {page_num}"

        # ── 1. Extract plain text ─────────────────────────────────────────
        text = page.get_text("text")
        if text.strip():
            text_by_source.append((source_label, text))
            all_text_parts.append(text)

        # ── 2. Extract tables (PyMuPDF 1.23+) ────────────────────────────
        try:
            table_finder = page.find_tables()
            for t_idx, table in enumerate(table_finder.tables, start=1):
                try:
                    table_data = table.extract()
                    if table_data:
                        table_text = _format_table(table_data, source_label)
                        if table_text.strip():
                            label = f"Table {t_idx} on {source_label}"
                            tables_by_source.append((label, table_text))
                            all_text_parts.append(table_text)
                except Exception:
                    pass
        except Exception:
            # Older PyMuPDF without find_tables
            pass

        # ── 3. Extract hyperlinks ─────────────────────────────────────────
        for link in page.get_links():
            uri = link.get("uri", "")
            if uri and uri.startswith("http"):
                links.append(uri)

    doc.close()

    full_text = "\n".join(all_text_parts)
    doc_type = _detect_doc_type(full_text)

    logger.info(
        f"PDF parsed: {len(text_by_source)} pages, "
        f"{len(tables_by_source)} tables, "
        f"{len(links)} links, type={doc_type}"
    )
    return {
        "text_by_source": text_by_source,
        "tables": tables_by_source,
        "links": list(set(links)),
        "doc_type": doc_type,
    }


# ─────────────────────────────────────────────────────────────────────────────
# DOCX Parser — with table extraction
# ─────────────────────────────────────────────────────────────────────────────
def parse_docx(file_path: str) -> Dict:
    """Extract text, tables, and URLs from a DOCX file."""
    try:
        from docx import Document
    except ImportError:
        raise RuntimeError("python-docx not installed: pip install python-docx")

    doc = Document(file_path)
    full_text_parts = []
    tables_by_source = []
    links: List[str] = []
    url_pattern = re.compile(r'https?://[^\s<>"{}|\\^`\[\]]+')

    # Paragraphs
    for para in doc.paragraphs:
        if para.text.strip():
            full_text_parts.append(para.text)

    # Tables
    for t_idx, table in enumerate(doc.tables, start=1):
        table_data = []
        for row in table.rows:
            row_cells = [cell.text.strip() for cell in row.cells]
            table_data.append(row_cells)
        table_text = _format_table(table_data, f"Table {t_idx}")
        if table_text.strip():
            label = f"Table {t_idx} in Document"
            tables_by_source.append((label, table_text))
            # NOTE: do NOT also append to full_text_parts. parse_document() merges
            # every table from `tables` into text_by_source as its own labeled
            # section, so adding it here too would embed the table content twice.

    # Hyperlinks from relationships
    for rel in doc.part.rels.values():
        if "hyperlink" in rel.reltype:
            url = rel.target_ref
            if url.startswith("http"):
                links.append(url)

    # Scan text for URLs
    for text in full_text_parts:
        links.extend(url_pattern.findall(text))

    combined_text = "\n".join(full_text_parts)
    doc_type = _detect_doc_type(combined_text)
    text_by_source = [("Document", combined_text)] if combined_text.strip() else []

    logger.info(f"DOCX parsed: {len(full_text_parts)} blocks, {len(tables_by_source)} tables, type={doc_type}")
    return {
        "text_by_source": text_by_source,
        "tables": tables_by_source,
        "links": list(set(links)),
        "doc_type": doc_type,
    }


# ─────────────────────────────────────────────────────────────────────────────
# Image Parser — OCR
# ─────────────────────────────────────────────────────────────────────────────
def parse_image(file_path: str) -> Dict:
    """Extract text from an image using Gemini's vision capability.
    No local OCR binary (Tesseract) required — works identically on a laptop and
    on the server, since it only needs the Gemini API key."""
    import io
    import time
    try:
        from PIL import Image
    except ImportError:
        raise RuntimeError("Pillow not installed")
    from google import genai
    from google.genai import types
    import config as cfg

    # Normalise any format (PNG/JPG/BMP/TIFF/…) to PNG bytes that Gemini accepts.
    img = Image.open(file_path)
    if img.mode not in ("RGB", "L"):
        img = img.convert("RGB")
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    png_bytes = buf.getvalue()

    client = genai.Client(
        api_key=cfg.GEMINI_API_KEY,
        http_options=types.HttpOptions(timeout=120_000),
    )
    prompt = (
        "Extract ALL text from this image exactly as it appears. Preserve numbers, "
        "dates, tables and the reading order, as plain text. Do not add any "
        "commentary or explanation. If the image contains no readable text, reply "
        "with nothing."
    )

    text = ""
    for attempt in range(4):
        try:
            resp = client.models.generate_content(
                model=cfg.CHAT_MODEL,
                contents=[
                    types.Part.from_bytes(data=png_bytes, mime_type="image/png"),
                    prompt,
                ],
            )
            text = (resp.text or "").strip()
            break
        except Exception as e:
            transient = any(s in str(e) for s in
                            ("503", "UNAVAILABLE", "429", "RESOURCE_EXHAUSTED", "500", "INTERNAL"))
            if transient and attempt < 3:
                time.sleep(min(5 * (2 ** attempt), 60))
            else:
                raise RuntimeError(f"Image text extraction failed: {e}")

    links: List[str] = []
    url_pattern = re.compile(r'https?://[^\s<>"{}|\\^`\[\]]+')
    links.extend(url_pattern.findall(text))

    doc_type = _detect_doc_type(text)
    text_by_source = [("Image", text)] if text.strip() else []
    logger.info(f"Image read via Gemini vision: {len(text.split())} words, type={doc_type}")
    return {
        "text_by_source": text_by_source,
        "tables": [],
        "links": list(set(links)),
        "doc_type": doc_type,
    }


# ─────────────────────────────────────────────────────────────────────────────
# PPTX Parser
# ─────────────────────────────────────────────────────────────────────────────
def parse_pptx(file_path: str) -> Dict:
    """Extract text and links from PowerPoint files."""
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("python-pptx not installed: pip install python-pptx")

    prs = Presentation(file_path)
    text_by_source = []
    links = []
    all_text = []
    url_pattern = re.compile(r'https?://[^\s<>"{}|\\^`\[\]]+')

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        if slide_text:
            combined = "\n".join(slide_text)
            text_by_source.append((f"Slide {slide_num}", combined))
            all_text.append(combined)
            links.extend(url_pattern.findall(combined))

    doc_type = _detect_doc_type("\n".join(all_text))
    return {
        "text_by_source": text_by_source,
        "tables": [],
        "links": list(set(links)),
        "doc_type": doc_type,
    }


# ─────────────────────────────────────────────────────────────────────────────
# Main dispatcher
# ─────────────────────────────────────────────────────────────────────────────
def parse_document(file_path: str) -> Dict:
    """Auto-detect file type and parse accordingly. Always includes tables."""
    ext = Path(file_path).suffix.lower()

    parsers = {
        ".pdf":  parse_pdf,
        ".docx": parse_docx,
        ".doc":  parse_docx,
        ".jpg":  parse_image,
        ".jpeg": parse_image,
        ".png":  parse_image,
        ".bmp":  parse_image,
        ".tiff": parse_image,
        ".tif":  parse_image,
        ".pptx": parse_pptx,
        ".ppt":  parse_pptx,
    }

    parser_fn = parsers.get(ext)
    if not parser_fn:
        raise ValueError(f"Unsupported file type: {ext}. Supported: {list(parsers.keys())}")

    result = parser_fn(file_path)

    # Merge tables into text_by_source as additional searchable content
    # (tables are labeled separately so FAISS can cite them)
    for label, table_text in result.get("tables", []):
        result["text_by_source"].append((label, table_text))

    # Scan all text for extra URLs — improved regex catches encoded chars, query params, fragments
    url_pattern = re.compile(
        r'https?://'                     # scheme
        r'[a-zA-Z0-9]'                   # must start with alphanumeric
        r'[^\s<>"\'{}|\\^`\[\]]*'        # match until whitespace or invalid chars
    )
    all_links = list(result.get("links", []))
    for _, text in result.get("text_by_source", []):
        raw_urls = url_pattern.findall(text)
        for url in raw_urls:
            # Clean trailing punctuation that often gets captured from prose
            url = url.rstrip(".,;:!?)>]}")
            if url and len(url) > 10:
                all_links.append(url)
    result["links"] = list(set(all_links))

    logger.info(f"Document parsed: {len(result['text_by_source'])} sections, "
                f"{len(result.get('tables', []))} tables, {len(result['links'])} unique links")

    return result
